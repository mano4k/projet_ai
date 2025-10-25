from flask import Flask, redirect, request, render_template, url_for, session, flash
from openai import OpenAI
from pathlib import Path
import os
import uuid
from langchain_openai import ChatOpenAI
from dotenv import load_dotenv
import openai   
from werkzeug.utils import secure_filename 

from pypdf import PdfReader #  PDF -> texte
from docx import Document as DocxDocument       # DOCX -> texte (paquet: python-docx)
from pptx import Presentation    

env_path = Path(__file__).resolve().parent / ".env"
print("Chemin absolu de .env :", env_path)
load_dotenv(dotenv_path=env_path)

# BASE_URL = "http://localhost:1234/v1"
#MODEL    = "meta-llama-3.1-8b-instruct"
# API_KEY  = "lm-studio"

#client= OpenAI(base_url="http://127.0.0.1:1234/v1", api_key="lm-studio")
# api_key = os.getenv("OPENAI_API_KEY")
# print("la cle:", api_key)
# if not api_key:
#     raise RuntimeError(f"Aucune clé API trouvée dans {env_path}")

load_dotenv()
OpenAI.api_key=os.getenv("OPENAI_API_KEY")

BASE_DIR = Path(__file__).resolve().parent 
app=Flask(__name__, template_folder=str(BASE_DIR / "templates"),  # dossier des templates Jinja2
    static_folder=str(BASE_DIR / "static"),       )
app.secret_key=os.getenv("SECRET_KEY") or "dev-secret-key"



###################################################

app.config["UPLOAD_FOLDER"] = str(BASE_DIR / "uploads")                 # binaires uploadés
app.config["TEXT_FOLDER"]   = str(BASE_DIR / "uploads" / "texts")       # textes extraits (.txt)
app.config["AUDIO_FOLDER"]  = str(BASE_DIR / "static"  / "audio")       # MP3 générés
app.config["MAX_CONTENT_LENGTH"] = 20 * 1024 * 1024  # 20 Mo max par upload

Path(app.config["UPLOAD_FOLDER"]).mkdir(parents=True, exist_ok=True)
Path(app.config["TEXT_FOLDER"]).mkdir(parents=True, exist_ok=True)
Path(app.config["AUDIO_FOLDER"]).mkdir(parents=True, exist_ok=True)

ALLOWED_EXT = {".pdf", ".docx", ".pptx"}

##################*****************************########################

def allowed_file(filename: str) -> bool:
    # True si l’extension est supportée. (Petite vérification sur le coté au cas oèu)
    return Path(filename).suffix.lower() in ALLOWED_EXT


def extract_text_from_pdf(filepath: str) -> str:
    # Extrait le texte d’un PDF (non scanné). Pour du scanné : prévoir un OCR.
    try:
        reader = PdfReader(filepath)
        pages = []
        for page in reader.pages:
            pages.append((page.extract_text() or "").strip())
        return "\n\n".join(pages).strip()
    except Exception as e:
        return f"[Erreur PDF] {e}"
    

def generate_summary(text: str, domaine: str = "générale", niveau_etude: str = "débutant", n_points: int = 5, max_chars: int = 8000, model: str = "gpt-5",) -> str:
    """
    Analyse un texte selon la discipline et le niveau d'étude :
    1) identifie les points essentiels sous forme de paragraphes,
    2) rédige un résumé clair et adapté au niveau d'étude.

    - discipline : domaine du texte (cybersécurité, biologie, économie…)
    - niveau_etude : débutant / intermédiaire / avancé / universitaire
    - n_points : nombre de points clés à générer
    - max_chars : limite de longueur du texte à analyser
    - model : modèle OpenAI utilisé
    """
    if not text:
        return "Aucun texte détecté dans le document."
    snippet = text[:max_chars]  
    prompt = f"""
Tu es un assistant expert en pédagogie et en vulgarisation scientifique.

Le texte suivant appartient à la **domaine**  :   {domaine}
L’audience cible est de **niveau d’étude** : {niveau_etude}.
Lis attentivement ce texte et effectue les étapes suivantes :
    **Points importants**  
    Identifie les {n_points} points ou idées les plus importants.  
    Rédige **chaque point comme un petit paragraphe** (3–5 phrases).  
    Commence chaque paragraphe par un **titre court en gras**, suivi d’une explication claire.  
    Utilise un langage adapté à un étudiant de niveau {niveau_etude}.  
    Ne fais pas de liste à puces : uniquement des paragraphes.
    
    **Résumé global**  
    Rédige ensuite un **résumé de synthèse** (250–350 mots maximum).  
    Adapte le ton, le vocabulaire et le niveau de détail à la discipline {domaine}.  
    Le résumé doit permettre à un étudiant de {niveau_etude} de comprendre les idées principales facilement.

    **Format de sortie Markdown attendu :**
 Points clés
[Paragraphe 1]
[Paragraphe 2]
[Paragraphe {n_points}]

# Résumé
[Paragraphe du résumé final]

--- TEXTE À ANALYSER ---
{snippet}
--- FIN DU TEXTE ---
""".strip()

    try:
        completion = openai.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
           
        )
        return completion.choices[0].message.content.strip()
    except Exception as e:
        return f"[Erreur OpenAI] {e}"
    
    
def extract_text_from_docx(filepath: str) -> str:
    # Extrait le texte d’un DOCX via python-docx.
    try:
        doc = DocxDocument(filepath)
        return "\n".join(p.text for p in doc.paragraphs).strip()
    except Exception as e:
        return f"[Erreur DOCX] {e}"

def extract_text_from_pptx(filepath: str) -> str:
    # Extrait le texte d’un PPTX via python-pptx (titres, zones, tableaux…).
    try:
        prs = Presentation(filepath)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chunks.append(shape.text.strip())
        return "\n\n".join(chunks).strip()
    except Exception as e:
        return f"[Erreur PPTX] {e}"
    


def extract_text_from_pdf(filepath: str) -> str:
    # Extrait le texte d’un PDF (non scanné). Pour du scanné : prévoir un OCR.
    try:
        reader = PdfReader(filepath)
        pages = []
        for page in reader.pages:
            pages.append((page.extract_text() or "").strip())
        return "\n\n".join(pages).strip()
    except Exception as e:
        return f"[Erreur PDF] {e}"
    

def extract_text_any(filepath: str) -> str:
    # Routeur simple par extension.
    ext = Path(filepath).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    if ext == ".docx":
        return extract_text_from_docx(filepath)
    if ext == ".pptx":
        return extract_text_from_pptx(filepath)
    return ""



def get_doc_text() -> str:
    
    p = session.get("doc_text_path")
    if not p:
        return ""
    try:
        if os.path.exists(p):
            return Path(p).read_text(encoding="utf-8")
    except Exception:
        pass
    return ""



def save_text_to_disk(source_path: str, text: str) -> str:
  
    stem = Path(source_path).stem  # nom du fichier sans extension
    txt_name = f"{stem}.txt"
    txt_path = os.path.join(app.config["TEXT_FOLDER"], txt_name)
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text or "")
    return txt_path

    
    
@app.route("/", methods=["GET", "POST"])
def index():
    # - POST : exécute l’action (upload/résumé/audio/qcm/soumission) puis redirect GET (PRG)
    # - GET  : affiche l’état courant.
    # La session ne contient que : niveau, filename, doc_text_path, summary, qcm, audio_url…
    
    # Valeurs par défaut (première visite)
    session.setdefault("niveau", "debutant")
    session.setdefault("filename", None)
    session.setdefault("domaine", "Informatique")
    session.setdefault("doc_text_path", None)  # <-- chemin du .txt (notre pivot)
    session.setdefault("summary", "")
  
    
    if request.method == "POST":
        action = request.form.get("action", "")
        # conserve le niveau (select latéral) ; défaut : débutant
        session["niveau"] = request.form.get("niveau") or session.get("niveau") or "debutant"
        session["domaine"] = request.form.get("domaine") or session.get("domaine") or "Informatique"
        niveau=session["niveau"],
        domain=session["domaine"],
        # (1) Upload d’un document (PDF/DOCX/PPTX)
        if action == "upload":
            f = request.files.get("file")
            if not f or f.filename == "":
                flash("Choisis un fichier PDF / DOCX / PPTX.", "error")
                return redirect(url_for("index"))
            if not allowed_file(f.filename):
                flash("Type de fichier non supporté.", "error")
                return redirect(url_for("index"))

            # Sauvegarde binaire
            safe_name = secure_filename(f.filename)
            save_name = f"{uuid.uuid4().hex}_{safe_name}"
            save_path = os.path.join(app.config["UPLOAD_FOLDER"], save_name)
            f.save(save_path)
            
            
            
             

            
            session["filename"] = save_name
            # Extraction du texte puis écriture sur disque .txt
            raw_text = extract_text_any(save_path)
            txt_path = save_text_to_disk(save_path, raw_text)

            # On stocke en session UNIQUEMENT les métadonnées / petits champs
            session["filename"] = os.path.basename(save_path)
            session["doc_text_path"] = txt_path
            session["summary"] = ""
           

            return redirect(url_for("index", msg="Fichier importé et texte extrait."))
        
        elif action == "resume":
            session["summary"] = generate_summary(get_doc_text(), domain,niveau)
           
            return redirect(url_for("index", msg="Résumé généré."))


    
    
    return render_template(
        "index.html",
        filename=session.get("filename"),
        summary=session.get("summary", ""),
        # qcm=session.get("qcm", []),
        # audio_url=session.get("audio_url", ""),
        niveau=session.get("niveau", "debutant"),
        domaine=session.get("domaine", "Informatique"),
        # results=session.get("results", []),
        score_line=session.get("score_line", ""),
    )

###################################################





if __name__ == "__main__":
    # En prod : debug=False et gunicorn derrière un reverse proxy.
    app.run(debug=True)
